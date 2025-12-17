#!/usr/bin/env python3
"""
Create both HTML slideshow and PowerPoint presentation from POWERPOINT_SLIDE_CONTENT.md
This script parses the markdown content and generates both formats.
"""

import re
from pathlib import Path
from typing import Dict, List, Optional

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError:
    print("ERROR: python-pptx not installed.")
    print("Please install it with: pip install python-pptx")
    exit(1)

try:
    from image_icons import get_font_awesome_class, get_icon_for_description
except ImportError:
    # Fallback if image_icons module not found
    def get_icon_for_description(description: str) -> str:
        desc_lower = description.lower()
        if "bronze" in desc_lower or "medal" in desc_lower:
            return "🥉"
        elif "silver" in desc_lower:
            return "🥈"
        elif "gold" in desc_lower:
            return "🥇"
        elif "clock" in desc_lower or "time" in desc_lower:
            return "⏰"
        elif "dollar" in desc_lower or "cost" in desc_lower or "money" in desc_lower:
            return "💰"
        elif "warning" in desc_lower or "alert" in desc_lower:
            return "⚠️"
        elif "checklist" in desc_lower or "list" in desc_lower:
            return "✅"
        elif "code" in desc_lower:
            return "💻"
        elif "database" in desc_lower or "table" in desc_lower:
            return "🗄️"
        elif (
            "dashboard" in desc_lower
            or "chart" in desc_lower
            or "analytics" in desc_lower
        ):
            return "📊"
        elif "cloud" in desc_lower or "databricks" in desc_lower:
            return "☁️"
        elif "logo" in desc_lower or "company" in desc_lower:
            return "🏢"
        return "📌"

    def get_font_awesome_class(description: str) -> str:
        desc_lower = description.lower()
        if "databricks" in desc_lower or "cloud" in desc_lower:
            return "fa-cloud"
        elif "database" in desc_lower or "table" in desc_lower:
            return "fa-database"
        elif "dashboard" in desc_lower or "analytics" in desc_lower:
            return "fa-chart-line"
        elif "code" in desc_lower:
            return "fa-code"
        elif "clock" in desc_lower or "time" in desc_lower:
            return "fa-clock"
        elif "dollar" in desc_lower or "money" in desc_lower:
            return "fa-dollar-sign"
        elif "warning" in desc_lower or "alert" in desc_lower:
            return "fa-exclamation-triangle"
        elif "checklist" in desc_lower or "list" in desc_lower:
            return "fa-tasks"
        return "fa-image"


# Color scheme matching the presentation
COLORS = {
    "bronze": RGBColor(205, 127, 50),  # #CD7F32
    "silver": RGBColor(192, 192, 192),  # #C0C0C0
    "gold": RGBColor(255, 215, 0),  # #FFD700
    "accent": RGBColor(0, 115, 230),  # Databricks blue
    "text": RGBColor(0, 0, 0),  # Black
    "background": RGBColor(255, 255, 255),  # White
    "code_bg": RGBColor(245, 245, 245),  # Light gray for code
}

COLOR_HEX = {
    "bronze": "#CD7F32",
    "silver": "#C0C0C0",
    "gold": "#FFD700",
    "accent": "#0073E6",
    "text": "#000000",
    "background": "#FFFFFF",
    "code_bg": "#F5F5F5",
}


class SlideContent:
    """Represents a single slide's content."""

    def __init__(self):
        self.title: Optional[str] = None
        self.subtitle: Optional[str] = None
        self.content: Optional[str] = None
        self.bullets: List[str] = []
        self.left_col: List[str] = []
        self.right_col: List[str] = []
        self.code: Optional[str] = None
        self.bottom_text: Optional[str] = None
        self.slide_type: str = "content"
        self.color: Optional[str] = None
        self.footer: Optional[str] = None
        self.quadrants: Dict[str, List[str]] = {}
        self.diagram: Optional[str] = None
        self.images: List[str] = []  # List of image descriptions/requirements


def parse_markdown_content(file_path: str) -> List[SlideContent]:
    """Parse POWERPOINT_SLIDE_CONTENT.md and extract all slide content."""
    with open(file_path, encoding="utf-8") as f:
        content = f.read()

    slides = []
    current_slide = None

    lines = content.split("\n")
    i = 0

    while i < len(lines):
        line = lines[i].strip()

        # Check for slide header: "### Slide X: Title"
        slide_match = re.match(r"^###\s+Slide\s+\d+[:\-]\s*(.+)$", line)
        if slide_match:
            if current_slide:
                # Ensure slide has at least a title before adding
                if not current_slide.title:
                    # Use the title from the previous header if available
                    pass
                slides.append(current_slide)
            current_slide = SlideContent()
            # Set initial title from header (may be overridden later)
            title_text = slide_match.group(1)
            current_slide.title = title_text
            # Try to infer slide type from title
            if "Title" in title_text:
                current_slide.slide_type = "title"
            elif "Agenda" in title_text:
                current_slide.slide_type = "bullets"
            elif (
                "Two-column" in title_text
                or "comparison" in title_text.lower()
                or "vs" in title_text.lower()
            ):
                current_slide.slide_type = "two-column"
            elif "Code" in title_text or "code" in title_text.lower():
                current_slide.slide_type = "code"
            elif "Diagram" in title_text:
                current_slide.slide_type = "content"
            elif "Four" in title_text or "Quadrant" in title_text:
                current_slide.slide_type = "four-quadrant"
            i += 1
            continue

        if not current_slide:
            i += 1
            continue

        # Parse different content types
        if line.startswith("**Title") and not current_slide.title:
            # Extract title from code block
            i += 1
            if i < len(lines) and lines[i].strip().startswith("```"):
                i += 1
                title_lines = []
                while i < len(lines) and not lines[i].strip().startswith("```"):
                    title_lines.append(lines[i])
                    i += 1
                title_text = "\n".join(title_lines).strip()
                if title_text:
                    current_slide.title = title_text

        elif line.startswith("**Subtitle"):
            i += 1
            if i < len(lines) and lines[i].strip().startswith("```"):
                i += 1
                subtitle_lines = []
                while i < len(lines) and not lines[i].strip().startswith("```"):
                    subtitle_lines.append(lines[i])
                    i += 1
                current_slide.subtitle = "\n".join(subtitle_lines).strip()

        elif line.startswith("**Content") or line.startswith("**Purpose"):
            i += 1
            if i < len(lines) and lines[i].strip().startswith("```"):
                i += 1
                content_lines = []
                while i < len(lines) and not lines[i].strip().startswith("```"):
                    content_lines.append(lines[i])
                    i += 1
                current_slide.content = "\n".join(content_lines).strip()

        elif "Bullet" in line or "Bullets" in line:
            i += 1
            if i < len(lines) and lines[i].strip().startswith("```"):
                i += 1
                bullet_lines = []
                while i < len(lines) and not lines[i].strip().startswith("```"):
                    bullet_lines.append(lines[i])
                    i += 1
                bullets_text = "\n".join(bullet_lines).strip()
                # Parse bullets (lines starting with •, -, or numbers)
                for bullet_line in bullets_text.split("\n"):
                    bullet_line = bullet_line.strip()
                    if bullet_line and (
                        bullet_line.startswith("•")
                        or bullet_line.startswith("-")
                        or bullet_line.startswith("✓")
                        or re.match(r"^\d+\.", bullet_line)
                    ):
                        # Remove bullet markers
                        clean_bullet = re.sub(r"^[•\-\✓]\s*", "", bullet_line)
                        clean_bullet = re.sub(r"^\d+\.\s*", "", clean_bullet)
                        if clean_bullet:
                            current_slide.bullets.append(clean_bullet)

        elif "Left Column" in line or ("Left Side" in line and "Sequential" in line):
            i += 1
            if i < len(lines) and lines[i].strip().startswith("```"):
                i += 1
                left_lines = []
                while i < len(lines) and not lines[i].strip().startswith("```"):
                    left_lines.append(lines[i])
                    i += 1
                left_text = "\n".join(left_lines).strip()
                # Check if this is code (starts with # or def or import)
                if (
                    left_text.strip().startswith("#")
                    or "def " in left_text
                    or "import " in left_text
                ):
                    # This is code, store as left_col but mark as code
                    current_slide.left_col = [left_text]  # Store full code block
                else:
                    # Regular bullets
                    for left_line in left_text.split("\n"):
                        left_line = left_line.strip()
                        if left_line:
                            # Handle lines that might not have bullet markers
                            if (
                                left_line.startswith("•")
                                or left_line.startswith("-")
                                or re.match(r"^\d+\.", left_line)
                            ):
                                clean_line = re.sub(r"^[•\-]\s*", "", left_line)
                                clean_line = re.sub(r"^\d+\.\s*", "", clean_line)
                            else:
                                clean_line = left_line
                            if clean_line and clean_line not in ["───────────"]:
                                current_slide.left_col.append(clean_line)

        elif "Right Column" in line or ("Right Side" in line and "Parallel" in line):
            i += 1
            if i < len(lines) and lines[i].strip().startswith("```"):
                i += 1
                right_lines = []
                while i < len(lines) and not lines[i].strip().startswith("```"):
                    right_lines.append(lines[i])
                    i += 1
                right_text = "\n".join(right_lines).strip()
                # Check if this is code (starts with # or def or import)
                if (
                    right_text.strip().startswith("#")
                    or "def " in right_text
                    or "import " in right_text
                ):
                    # This is code, store as right_col but mark as code
                    current_slide.right_col = [right_text]  # Store full code block
                else:
                    # Regular bullets
                    for right_line in right_text.split("\n"):
                        right_line = right_line.strip()
                        if right_line:
                            # Handle lines that might not have bullet markers
                            if (
                                right_line.startswith("•")
                                or right_line.startswith("-")
                                or re.match(r"^\d+\.", right_line)
                            ):
                                clean_line = re.sub(r"^[•\-]\s*", "", right_line)
                                clean_line = re.sub(r"^\d+\.\s*", "", clean_line)
                            else:
                                clean_line = right_line
                            if clean_line and clean_line not in ["───────────"]:
                                current_slide.right_col.append(clean_line)

        elif "Code Example" in line or ("Code" in line and "Code:" not in line):
            i += 1
            # Look for code block (```python, ```json, or just ```)
            while i < len(lines):
                if lines[i].strip().startswith("```"):
                    # Found code block start
                    i += 1
                    code_lines = []
                    while i < len(lines) and not lines[i].strip().startswith("```"):
                        code_lines.append(lines[i])
                        i += 1
                    if code_lines:
                        current_slide.code = "\n".join(code_lines).strip()
                    break
                i += 1

        elif "Bottom Text" in line:
            i += 1
            if i < len(lines) and lines[i].strip().startswith("```"):
                i += 1
                bottom_lines = []
                while i < len(lines) and not lines[i].strip().startswith("```"):
                    bottom_lines.append(lines[i])
                    i += 1
                current_slide.bottom_text = "\n".join(bottom_lines).strip()

        elif "Footer Text" in line:
            i += 1
            if i < len(lines) and lines[i].strip().startswith("```"):
                i += 1
                footer_lines = []
                while i < len(lines) and not lines[i].strip().startswith("```"):
                    footer_lines.append(lines[i])
                    i += 1
                current_slide.footer = "\n".join(footer_lines).strip()

        elif "Color:" in line:
            if "bronze" in line.lower():
                current_slide.color = "bronze"
            elif "silver" in line.lower():
                current_slide.color = "silver"
            elif "gold" in line.lower():
                current_slide.color = "gold"

        elif "Diagram" in line:
            i += 1
            if i < len(lines) and lines[i].strip().startswith("```"):
                i += 1
                diagram_lines = []
                while i < len(lines) and not lines[i].strip().startswith("```"):
                    diagram_lines.append(lines[i])
                    i += 1
                current_slide.diagram = "\n".join(diagram_lines).strip()

        elif "Images to Insert" in line or "Images:" in line:
            i += 1
            image_lines = []
            # Collect image descriptions until next section
            while i < len(lines):
                line_check = lines[i].strip()
                if (
                    line_check.startswith("**")
                    or line_check.startswith("---")
                    or line_check.startswith("###")
                ):
                    break
                if line_check.startswith("-"):
                    image_lines.append(line_check[1:].strip())
                i += 1
            current_slide.images = image_lines
            continue

        elif (
            "Top Left" in line
            or "Top Right" in line
            or "Bottom Left" in line
            or "Bottom Right" in line
        ):
            quadrant = line.split("-")[0].strip().replace("**", "").lower()
            i += 1
            if i < len(lines) and lines[i].strip().startswith("```"):
                i += 1
                quadrant_lines = []
                while i < len(lines) and not lines[i].strip().startswith("```"):
                    quadrant_lines.append(lines[i])
                    i += 1
                quadrant_text = "\n".join(quadrant_lines).strip()
                bullets = []
                for q_line in quadrant_text.split("\n"):
                    q_line = q_line.strip()
                    if q_line and (q_line.startswith("✓") or q_line.startswith("•")):
                        clean_line = re.sub(r"^[✓•]\s*", "", q_line)
                        if clean_line:
                            bullets.append(clean_line)
                if quadrant not in current_slide.quadrants:
                    current_slide.quadrants[quadrant] = []
                current_slide.quadrants[quadrant] = bullets

        # Infer slide type from content
        if current_slide.left_col and current_slide.right_col:
            current_slide.slide_type = "two-column"
        elif current_slide.bullets and not current_slide.code:
            current_slide.slide_type = "bullets"
        elif current_slide.code:
            current_slide.slide_type = "code"
        elif current_slide.quadrants:
            current_slide.slide_type = "four-quadrant"

        i += 1

    # Add last slide
    if current_slide:
        slides.append(current_slide)

    return slides


def create_html_slideshow(
    slides: List[SlideContent], output_file: str = "presentation_slideshow.html"
):
    """Create HTML slideshow using Reveal.js framework."""

    html_content = """<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="utf-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Pipeline Builder on Databricks</title>

    <!-- Reveal.js CSS -->
    <link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/reveal.js@4.3.1/dist/reveal.css">
    <link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/reveal.js@4.3.1/dist/theme/white.css">

    <!-- Highlight.js for code syntax highlighting -->
    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/highlight.js/11.9.0/styles/default.min.css">
    <script src="https://cdnjs.cloudflare.com/ajax/libs/highlight.js/11.9.0/highlight.min.js"></script>
    <script src="https://cdnjs.cloudflare.com/ajax/libs/highlight.js/11.9.0/languages/python.min.js"></script>

    <!-- Font Awesome for icons -->
    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/6.4.0/css/all.min.css">

    <style>
        :root {
            --bronze: #CD7F32;
            --silver: #C0C0C0;
            --gold: #FFD700;
            --accent: #0066CC;
            --accent-dark: #004499;
            --accent-light: #3399FF;
            --text: #1A1A1A;
            --text-light: #4A4A4A;
            --text-muted: #6B6B6B;
            --background: #FFFFFF;
            --background-alt: #F8F9FA;
            --background-subtle: #FAFBFC;
            --code-bg: #1E1E1E;
            --code-text: #D4D4D4;
            --border: #E1E8ED;
            --gradient-primary: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            --gradient-accent: linear-gradient(135deg, #0066CC 0%, #004499 100%);
            --gradient-subtle: linear-gradient(180deg, #FFFFFF 0%, #F8F9FA 100%);
            --shadow: 0 2px 8px rgba(0, 0, 0, 0.08);
            --shadow-md: 0 4px 12px rgba(0, 0, 0, 0.1);
            --shadow-lg: 0 8px 24px rgba(0, 0, 0, 0.12);
        }

        * {
            box-sizing: border-box;
        }

        .reveal {
            font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', 'Roboto', 'Helvetica Neue', Arial, sans-serif;
            background: linear-gradient(135deg, #F5F7FA 0%, #E8ECF1 50%, #DDE4EA 100%);
            color: var(--text);
        }

        .reveal .slides section {
            text-align: left;
            background: var(--background);
            border-radius: 16px;
            padding: 3em 2.5em;
            box-shadow: var(--shadow-lg);
            border: 1px solid var(--border);
            position: relative;
            overflow: hidden;
        }

        .reveal .slides section::before {
            content: '';
            position: absolute;
            top: 0;
            left: 0;
            right: 0;
            height: 4px;
            background: var(--gradient-accent);
        }

        .reveal h1, .reveal h2, .reveal h3 {
            color: var(--text);
            font-weight: 700;
            letter-spacing: -0.02em;
            line-height: 1.2;
        }

        .reveal h1 {
            font-size: 2.8em;
            margin-bottom: 0.6em;
            margin-top: 0.2em;
            color: var(--accent);
            font-weight: 800;
            border-bottom: 3px solid var(--accent);
            padding-bottom: 0.3em;
            display: inline-block;
            width: 100%;
        }

        .reveal h2 {
            font-size: 1.9em;
            margin-bottom: 0.5em;
            color: var(--accent-dark);
            font-weight: 700;
        }

        .reveal h3 {
            font-size: 1.4em;
            color: var(--text-light);
            font-weight: 600;
        }

        .reveal .subtitle {
            font-size: 1.6em;
            font-weight: 600;
            color: var(--accent);
            margin-top: 0.4em;
            margin-bottom: 0.8em;
            padding: 0.4em 0;
            border-left: 4px solid var(--accent);
            padding-left: 0.8em;
        }

        .reveal ul, .reveal ol {
            margin-left: 1.2em;
            font-size: 0.95em;
            line-height: 1.8;
        }

        .reveal li {
            margin-bottom: 0.7em;
            color: var(--text-light);
            line-height: 1.7;
            padding-left: 0.5em;
        }

        .reveal li::marker {
            color: var(--accent);
            font-weight: bold;
        }

        .reveal p {
            color: var(--text-light);
            line-height: 1.7;
            font-size: 0.95em;
            margin-bottom: 1em;
        }

        .reveal .two-column {
            display: grid;
            grid-template-columns: 1fr 1fr;
            gap: 2.5em;
            margin-top: 1.5em;
        }

        .reveal .two-column > div {
            background: var(--background-subtle);
            padding: 2em 1.8em;
            border-radius: 12px;
            border-left: 5px solid var(--accent);
            box-shadow: var(--shadow-md);
            transition: transform 0.2s, box-shadow 0.2s;
        }

        .reveal .two-column > div:hover {
            transform: translateY(-2px);
            box-shadow: var(--shadow-lg);
        }

        .reveal .four-quadrant {
            display: grid;
            grid-template-columns: 1fr 1fr;
            grid-template-rows: 1fr 1fr;
            gap: 1.5em;
        }

        .reveal .quadrant {
            padding: 1.5em;
            border: 2px solid var(--accent);
            border-radius: 12px;
            background: var(--background-alt);
            box-shadow: var(--shadow);
            transition: transform 0.2s;
        }

        .reveal .quadrant:hover {
            transform: translateY(-2px);
            box-shadow: var(--shadow-lg);
        }

        .reveal .code-block {
            background: var(--code-bg);
            border: none;
            border-radius: 12px;
            padding: 1.8em;
            font-family: 'SF Mono', 'Monaco', 'Consolas', 'Courier New', monospace;
            font-size: 0.72em;
            overflow-x: auto;
            text-align: left;
            box-shadow: var(--shadow-lg);
            position: relative;
            margin: 1.5em 0;
        }

        .reveal .code-block::before {
            content: '';
            position: absolute;
            top: 0;
            left: 0;
            right: 0;
            height: 4px;
            background: var(--gradient-accent);
            border-radius: 12px 12px 0 0;
        }

        .reveal .code-block::after {
            content: 'Python';
            position: absolute;
            top: 0.8em;
            right: 1.2em;
            font-size: 0.7em;
            color: rgba(255, 255, 255, 0.5);
            font-family: 'SF Mono', 'Monaco', 'Consolas', monospace;
            text-transform: uppercase;
            letter-spacing: 0.1em;
        }

        .reveal .code-block pre {
            margin: 0;
            color: var(--code-text);
        }

        .reveal .code-block code {
            background: transparent;
            padding: 0;
            color: var(--code-text);
        }

        .reveal .bottom-text {
            text-align: center;
            font-size: 1.4em;
            font-weight: 700;
            margin-top: 1.5em;
            padding: 1em 1.5em;
            background: var(--gradient-accent);
            color: white;
            border-radius: 12px;
            box-shadow: var(--shadow-md);
            letter-spacing: 0.02em;
        }

        .reveal .bronze-accent {
            color: var(--bronze);
            font-weight: 600;
        }

        .reveal .silver-accent {
            color: var(--silver);
            font-weight: 600;
        }

        .reveal .gold-accent {
            color: var(--gold);
            font-weight: 600;
        }

        .reveal .title-slide {
            text-align: center;
            background: var(--gradient-primary);
            color: white;
            border-radius: 16px;
            padding: 4em 3em;
            position: relative;
            overflow: hidden;
        }

        .reveal .title-slide::before {
            content: '';
            position: absolute;
            top: -50%;
            right: -50%;
            width: 200%;
            height: 200%;
            background: radial-gradient(circle, rgba(255, 255, 255, 0.1) 0%, transparent 70%);
            animation: pulse 20s ease-in-out infinite;
        }

        @keyframes pulse {
            0%, 100% { transform: scale(1) rotate(0deg); }
            50% { transform: scale(1.1) rotate(180deg); }
        }

        .reveal .title-slide h1 {
            color: white;
            -webkit-text-fill-color: white;
            text-shadow: 0 4px 20px rgba(0, 0, 0, 0.3);
            margin-top: 0.5em;
            border: none;
            font-size: 3.2em;
            font-weight: 900;
            letter-spacing: -0.03em;
            position: relative;
            z-index: 1;
        }

        .reveal .title-slide .subtitle {
            color: rgba(255, 255, 255, 0.95);
            -webkit-text-fill-color: rgba(255, 255, 255, 0.95);
            text-shadow: 0 2px 10px rgba(0, 0, 0, 0.2);
            border: none;
            padding: 0;
            margin-top: 0.8em;
            font-size: 1.8em;
            position: relative;
            z-index: 1;
        }

        .reveal .footer {
            font-size: 0.65em;
            color: rgba(255, 255, 255, 0.85);
            margin-top: 3em;
            position: relative;
            z-index: 1;
            font-weight: 500;
        }

        .reveal .diagram {
            font-family: 'SF Mono', 'Monaco', 'Courier New', monospace;
            font-size: 0.85em;
            white-space: pre;
            background: var(--background-subtle);
            padding: 2em;
            border-radius: 12px;
            border: 2px solid var(--border);
            box-shadow: var(--shadow-md);
            color: var(--text-light);
            line-height: 1.6;
        }

        .reveal .progress {
            background: var(--gradient-accent);
            height: 4px;
        }

        .reveal .controls {
            color: var(--accent);
        }

        .reveal .controls button {
            opacity: 0.8;
            transition: opacity 0.2s;
        }

        .reveal .controls button:hover {
            opacity: 1;
        }

        /* Slide transitions */
        .reveal .slides section {
            transition: all 0.3s ease;
        }

        /* Print styles */
        @media print {
            .reveal .slides section {
                page-break-after: always;
                box-shadow: none;
                border: 1px solid #ddd;
            }
        }

        /* Image and icon styles */
        .slide-icon {
            font-size: 2em;
            margin: 0.3em;
            display: inline-block;
            vertical-align: middle;
        }

        .slide-icon.fa {
            color: var(--accent);
            margin-right: 0.5em;
        }

        .image-placeholder {
            background: var(--background-alt);
            border: 2px dashed var(--border);
            border-radius: 8px;
            padding: 1em;
            margin: 0.5em 0;
            text-align: center;
            color: var(--text-muted);
            font-size: 0.85em;
        }

        .title-slide .slide-icon {
            color: rgba(255, 255, 255, 0.9);
            font-size: 3em;
        }

        .icon-grid {
            display: flex;
            flex-wrap: wrap;
            gap: 1em;
            margin: 1em 0;
            justify-content: center;
        }
    </style>
</head>
<body>
    <div class="reveal">
        <div class="slides">
"""

    for slide in slides:
        html_content += "            <section>\n"

        if slide.slide_type == "title":
            html_content += '                <div class="title-slide">\n'
            # Add icons for title slide
            if slide.images:
                html_content += '                    <div class="icon-grid">\n'
                for img_desc in slide.images[:3]:  # Limit to 3 icons
                    icon = get_icon_for_description(img_desc)
                    html_content += f'                        <span class="slide-icon">{icon}</span>\n'
                html_content += "                    </div>\n"
            if slide.title:
                html_content += (
                    f"                    <h1>{escape_html(slide.title)}</h1>\n"
                )
            if slide.subtitle:
                html_content += f'                    <h2 class="subtitle">{escape_html(slide.subtitle)}</h2>\n'
            if slide.footer:
                html_content += f'                    <div class="footer">{escape_html(slide.footer)}</div>\n'
            html_content += "                </div>\n"

        elif slide.slide_type == "two-column":
            if slide.title:
                html_content += f"                <h1>{escape_html(slide.title)}</h1>\n"
            # Add icons if specified
            if slide.images:
                html_content += '                <div style="margin: 0.5em 0 1em 0; text-align: center;">\n'
                for img_desc in slide.images[:3]:
                    icon = get_icon_for_description(img_desc)
                    fa_class = get_font_awesome_class(img_desc)
                    html_content += f'                    <i class="fas {fa_class} slide-icon" title="{escape_html(img_desc)}"></i>\n'
                html_content += "                </div>\n"
            html_content += '                <div class="two-column">\n'
            if slide.left_col:
                html_content += "                    <div>\n"
                # Check if this is code (code comparison slide)
                if len(slide.left_col) == 1 and (
                    "def " in slide.left_col[0]
                    or "import " in slide.left_col[0]
                    or slide.left_col[0].strip().startswith("#")
                ):
                    html_content += '                        <div class="code-block">\n'
                    html_content += f'                            <pre><code class="language-python">{escape_html(slide.left_col[0])}</code></pre>\n'
                    html_content += "                        </div>\n"
                else:
                    html_content += "                        <ul>\n"
                    for item in slide.left_col:
                        html_content += f"                            <li>{escape_html(item)}</li>\n"
                    html_content += "                        </ul>\n"
                html_content += "                    </div>\n"
            if slide.right_col:
                html_content += "                    <div>\n"
                # Check if this is code (code comparison slide)
                if len(slide.right_col) == 1 and (
                    "def " in slide.right_col[0]
                    or "import " in slide.right_col[0]
                    or slide.right_col[0].strip().startswith("#")
                ):
                    html_content += '                        <div class="code-block">\n'
                    html_content += f'                            <pre><code class="language-python">{escape_html(slide.right_col[0])}</code></pre>\n'
                    html_content += "                        </div>\n"
                else:
                    html_content += "                        <ul>\n"
                    for item in slide.right_col:
                        html_content += f"                            <li>{escape_html(item)}</li>\n"
                    html_content += "                        </ul>\n"
                html_content += "                    </div>\n"
            html_content += "                </div>\n"
            if slide.bottom_text:
                html_content += f'                <div class="bottom-text">{escape_html(slide.bottom_text)}</div>\n'

        elif slide.slide_type == "four-quadrant":
            if slide.title:
                html_content += f"                <h1>{escape_html(slide.title)}</h1>\n"
            html_content += '                <div class="four-quadrant">\n'
            for quadrant in ["top left", "top right", "bottom left", "bottom right"]:
                if quadrant in slide.quadrants:
                    html_content += '                    <div class="quadrant">\n'
                    html_content += "                        <ul>\n"
                    for item in slide.quadrants[quadrant]:
                        html_content += f"                            <li>{escape_html(item)}</li>\n"
                    html_content += "                        </ul>\n"
                    html_content += "                    </div>\n"
            html_content += "                </div>\n"

        elif slide.slide_type == "code":
            if slide.title:
                html_content += f"                <h1>{escape_html(slide.title)}</h1>\n"
            if slide.subtitle:
                html_content += f'                <h2 class="subtitle">{escape_html(slide.subtitle)}</h2>\n'
            if slide.code:
                html_content += '                <div class="code-block">\n'
                html_content += f'                    <pre><code class="language-python">{escape_html(slide.code)}</code></pre>\n'
                html_content += "                </div>\n"
            if slide.bullets:
                html_content += "                <ul>\n"
                for bullet in slide.bullets:
                    html_content += (
                        f"                    <li>{escape_html(bullet)}</li>\n"
                    )
                html_content += "                </ul>\n"
            if slide.content:
                html_content += f"                <p>{escape_html(slide.content)}</p>\n"

        elif slide.slide_type == "bullets":
            if slide.title:
                html_content += f"                <h1>{escape_html(slide.title)}</h1>\n"
            # Add icons if specified
            if slide.images:
                html_content += '                <div style="margin: 0.5em 0 1em 0;">\n'
                for img_desc in slide.images[:4]:  # Limit to 4 icons
                    icon = get_icon_for_description(img_desc)
                    fa_class = get_font_awesome_class(img_desc)
                    html_content += f'                    <i class="fas {fa_class} slide-icon" title="{escape_html(img_desc)}"></i>\n'
                html_content += "                </div>\n"
            if slide.subtitle:
                html_content += f'                <h2 class="subtitle">{escape_html(slide.subtitle)}</h2>\n'
            if slide.bullets:
                html_content += "                <ul>\n"
                for bullet in slide.bullets:
                    html_content += (
                        f"                    <li>{escape_html(bullet)}</li>\n"
                    )
                html_content += "                </ul>\n"

        else:  # content
            if slide.title:
                color_class = ""
                if slide.color:
                    color_class = f' class="{slide.color}-accent"'
                html_content += f"                <h1{color_class}>{escape_html(slide.title)}</h1>\n"
            # Add icons if specified (especially for medal layers)
            if slide.images:
                html_content += '                <div style="margin: 0.5em 0 1em 0; text-align: center;">\n'
                for img_desc in slide.images[:4]:
                    icon = get_icon_for_description(img_desc)
                    fa_class = get_font_awesome_class(img_desc)
                    # Use emoji for medals, Font Awesome for others
                    if (
                        "medal" in img_desc.lower()
                        or "bronze" in img_desc.lower()
                        or "silver" in img_desc.lower()
                        or "gold" in img_desc.lower()
                    ):
                        html_content += f'                    <span class="slide-icon" style="font-size: 3em;">{icon}</span>\n'
                    else:
                        html_content += f'                    <i class="fas {fa_class} slide-icon" title="{escape_html(img_desc)}"></i>\n'
                html_content += "                </div>\n"
            if slide.subtitle:
                html_content += f'                <h2 class="subtitle">{escape_html(slide.subtitle)}</h2>\n'
            if slide.content:
                html_content += f"                <p>{escape_html(slide.content)}</p>\n"
            if slide.diagram:
                html_content += f'                <div class="diagram">{escape_html(slide.diagram)}</div>\n'
            if slide.bullets:
                html_content += "                <ul>\n"
                for bullet in slide.bullets:
                    html_content += (
                        f"                    <li>{escape_html(bullet)}</li>\n"
                    )
                html_content += "                </ul>\n"

        html_content += "            </section>\n"

    html_content += """        </div>
    </div>

    <!-- Reveal.js JavaScript -->
    <script src="https://cdn.jsdelivr.net/npm/reveal.js@4.3.1/dist/reveal.js"></script>
    <script>
        Reveal.initialize({
            hash: true,
            controls: true,
            progress: true,
            center: true,
            transition: 'slide'
        });

        // Initialize syntax highlighting
        hljs.highlightAll();
    </script>
</body>
</html>"""

    with open(output_file, "w", encoding="utf-8") as f:
        f.write(html_content)

    print(f"✅ HTML slideshow created: {output_file}")
    return output_file


def escape_html(text: str) -> str:
    """Escape HTML special characters."""
    if not text:
        return ""
    return (
        text.replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
        .replace("'", "&#39;")
    )


def create_powerpoint(
    slides: List[SlideContent],
    output_file: str = "Pipeline_Builder_Databricks_Presentation_HTML.pptx",
):
    """Create PowerPoint presentation from slide content."""
    print("Creating PowerPoint presentation...")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print(f"Creating {len(slides)} slides...")

    for i, slide in enumerate(slides, 1):
        print(f"  [{i}/{len(slides)}] {slide.title or 'Untitled'}")
        create_pptx_slide(prs, slide)

    print(f"\nSaving presentation to: {output_file}")
    prs.save(output_file)
    print("✅ Presentation created successfully!")
    print(f"   File: {output_file}")
    print(f"   Total slides: {len(slides)}")

    return output_file


def create_pptx_slide(prs: Presentation, slide: SlideContent):
    """Create a single PowerPoint slide."""
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    pptx_slide = prs.slides.add_slide(blank_slide_layout)

    # Add subtle background color (except for title slides which get special treatment)
    if slide.slide_type != "title":
        background = pptx_slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(250, 251, 252)  # Very light gray background

    if slide.slide_type == "title":
        create_title_slide(pptx_slide, slide)
    elif slide.slide_type == "bullets":
        create_bullet_slide(pptx_slide, slide)
    elif slide.slide_type == "two-column":
        create_two_column_slide(pptx_slide, slide)
    elif slide.slide_type == "code":
        create_code_slide(pptx_slide, slide)
    elif slide.slide_type == "four-quadrant":
        create_four_quadrant_slide(pptx_slide, slide)
    else:
        create_content_slide(pptx_slide, slide)


def create_title_slide(slide, data: SlideContent):
    """Create title slide with gradient background."""
    # Add gradient-like background rectangle
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(102, 126, 234)  # Purple-blue gradient start
    bg_shape.line.fill.background()  # No border

    if data.title:
        add_title(slide, data.title, font_size=44, top=Inches(2))
        # Make title white on gradient background
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame.paragraphs:
                if shape.text_frame.paragraphs[0].text == data.title:
                    shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(
                        255, 255, 255
                    )
    if data.subtitle:
        add_text(slide, data.subtitle, font_size=32, top=Inches(3.5), bold=True)
        # Make subtitle white
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame.paragraphs:
                if shape.text_frame.paragraphs[0].text == data.subtitle:
                    shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(
                        255, 255, 255
                    )
    if data.footer:
        add_text(slide, data.footer, font_size=18, top=Inches(6.5))
        # Make footer light
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame.paragraphs:
                if shape.text_frame.paragraphs[0].text == data.footer:
                    shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(
                        240, 240, 240
                    )


def create_bullet_slide(slide, data: SlideContent):
    """Create slide with bullet points."""
    if data.title:
        add_title(slide, data.title, font_size=44)
    if data.subtitle:
        add_text(slide, data.subtitle, font_size=32, top=Inches(1.2), bold=True)
    if data.bullets:
        top_pos = Inches(1.8 if data.subtitle else 1.5)
        add_bullets(slide, data.bullets, top=top_pos)
    if data.code:
        add_code_block(slide, data.code, top=Inches(3))


def create_two_column_slide(slide, data: SlideContent):
    """Create two-column slide."""
    if data.title:
        add_title(slide, data.title, font_size=44)

    if data.left_col:
        # Check if this is code (code comparison slide)
        if len(data.left_col) == 1 and (
            "def " in data.left_col[0]
            or "import " in data.left_col[0]
            or data.left_col[0].strip().startswith("#")
        ):
            add_code_block(
                slide,
                data.left_col[0],
                left=Inches(0.5),
                top=Inches(1.5),
                width=Inches(4.25),
                height=Inches(5.0),
                font_size=14,
            )
        else:
            add_bullets(
                slide,
                data.left_col,
                left=Inches(0.5),
                top=Inches(1.5),
                width=Inches(4.25),
                font_size=20,
            )

    if data.right_col:
        # Check if this is code (code comparison slide)
        if len(data.right_col) == 1 and (
            "def " in data.right_col[0]
            or "import " in data.right_col[0]
            or data.right_col[0].strip().startswith("#")
        ):
            add_code_block(
                slide,
                data.right_col[0],
                left=Inches(5.25),
                top=Inches(1.5),
                width=Inches(4.25),
                height=Inches(5.0),
                font_size=14,
            )
        else:
            add_bullets(
                slide,
                data.right_col,
                left=Inches(5.25),
                top=Inches(1.5),
                width=Inches(4.25),
                font_size=20,
            )

    if data.bottom_text:
        add_text(
            slide,
            data.bottom_text,
            font_size=28,
            top=Inches(6.5),
            bold=True,
            center=True,
        )


def create_code_slide(slide, data: SlideContent):
    """Create slide with code."""
    if data.title:
        add_title(slide, data.title, font_size=44)
    if data.subtitle:
        add_text(slide, data.subtitle, font_size=32, top=Inches(1.2), bold=True)
    if data.code:
        top_pos = Inches(1.5 if not data.subtitle else 2.0)
        add_code_block(slide, data.code, top=top_pos)
    if data.bullets:
        add_bullets(slide, data.bullets, top=Inches(5), font_size=18)
    if data.content:
        add_text(slide, data.content, top=Inches(5.5), font_size=18)


def create_four_quadrant_slide(slide, data: SlideContent):
    """Create four-quadrant slide."""
    if data.title:
        add_title(slide, data.title, font_size=44)

    # Top left
    if "top left" in data.quadrants:
        add_bullets(
            slide,
            data.quadrants["top left"],
            left=Inches(0.5),
            top=Inches(1.8),
            width=Inches(4.25),
            font_size=18,
        )

    # Top right
    if "top right" in data.quadrants:
        add_bullets(
            slide,
            data.quadrants["top right"],
            left=Inches(5.25),
            top=Inches(1.8),
            width=Inches(4.25),
            font_size=18,
        )

    # Bottom left
    if "bottom left" in data.quadrants:
        add_bullets(
            slide,
            data.quadrants["bottom left"],
            left=Inches(0.5),
            top=Inches(4.0),
            width=Inches(4.25),
            font_size=18,
        )

    # Bottom right
    if "bottom right" in data.quadrants:
        add_bullets(
            slide,
            data.quadrants["bottom right"],
            left=Inches(5.25),
            top=Inches(4.0),
            width=Inches(4.25),
            font_size=18,
        )


def create_content_slide(slide, data: SlideContent):
    """Create content slide."""
    if data.title:
        add_title(slide, data.title, font_size=44)
    if data.subtitle:
        add_text(slide, data.subtitle, font_size=32, top=Inches(1.2), bold=True)
    if data.content:
        top_pos = Inches(1.8 if data.subtitle else 1.5)
        add_text(slide, data.content, top=top_pos)
    if data.diagram:
        add_text(slide, data.diagram, top=Inches(2.5), font_size=14)
    if data.bullets:
        top_pos = Inches(2.0 if data.content else 1.8)
        add_bullets(slide, data.bullets, top=top_pos, font_size=20)


def add_title(slide, text: str, font_size: int = 44, top: float = Inches(0.3)):
    """Add title to slide with professional styling."""
    left = Inches(0.5)
    width = Inches(9)
    height = Inches(1.2)

    # Add accent line under title
    line_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top + Inches(0.85), Inches(2), Inches(0.08)
    )
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = COLORS["accent"]
    line_shape.line.fill.background()

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP

    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.name = "Calibri"
    p.font.color.rgb = COLORS["accent"]
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(8)


def add_text(
    slide,
    text: str,
    font_size: int = 24,
    top: float = Inches(1.5),
    left: float = Inches(0.5),
    width: float = Inches(9),
    bold: bool = False,
    center: bool = False,
):
    """Add text to slide."""
    height = Inches(1)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP

    lines = text.split("\n")
    for i, line in enumerate(lines):
        if line.strip():
            p = (
                text_frame.paragraphs[i]
                if i < len(text_frame.paragraphs)
                else text_frame.add_paragraph()
            )
            p.text = line.strip()
            p.font.size = Pt(font_size)
            p.font.bold = bold
            p.font.color.rgb = COLORS["text"]
            p.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT
            p.space_after = Pt(4)


def add_bullets(
    slide,
    bullets: List[str],
    font_size: int = 24,
    top: float = Inches(1.5),
    left: float = Inches(0.5),
    width: float = Inches(9),
    height: float = Inches(5.5),
):
    """Add bullet points to slide."""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP

    for i, bullet in enumerate(bullets):
        if bullet.strip():
            p = (
                text_frame.paragraphs[i]
                if i < len(text_frame.paragraphs)
                else text_frame.add_paragraph()
            )
            p.text = bullet.strip()
            p.font.size = Pt(font_size)
            p.font.color.rgb = COLORS["text"]
            p.level = 0
            p.space_after = Pt(6)


def add_code_block(
    slide,
    code: str,
    font_size: int = 18,
    top: float = Inches(1.5),
    left: float = None,
    width: float = None,
    height: float = None,
):
    """Add code block to slide."""
    if left is None:
        left = Inches(0.5)
    if width is None:
        width = Inches(9)
    if height is None:
        height = Inches(5.5)

    # Add background rectangle with accent border
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(30, 30, 30)  # Dark background for code
    shape.line.color.rgb = COLORS["accent"]  # Accent color border
    shape.line.width = Pt(3)  # Thicker border

    # Add text box on top
    textbox = slide.shapes.add_textbox(
        left + Inches(0.1), top + Inches(0.1), width - Inches(0.2), height - Inches(0.2)
    )
    text_frame = textbox.text_frame
    text_frame.word_wrap = False
    text_frame.vertical_anchor = MSO_ANCHOR.TOP

    # Use monospace font for code
    lines = code.split("\n")
    for i, line in enumerate(lines[:35]):  # Limit to 35 lines
        p = (
            text_frame.paragraphs[i]
            if i < len(text_frame.paragraphs)
            else text_frame.add_paragraph()
        )
        p.text = line
        p.font.name = "Consolas"
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(212, 212, 212)  # Light text on dark background
        p.space_after = Pt(2)


def main():
    """Main function to generate both HTML and PowerPoint."""
    content_file = "POWERPOINT_SLIDE_CONTENT.md"

    if not Path(content_file).exists():
        print(f"ERROR: {content_file} not found!")
        return

    print(f"Parsing {content_file}...")
    slides = parse_markdown_content(content_file)
    print(f"✅ Parsed {len(slides)} slides")

    print("\n" + "=" * 60)
    print("Generating HTML slideshow...")
    print("=" * 60)
    create_html_slideshow(slides)

    print("\n" + "=" * 60)
    print("Generating PowerPoint presentation...")
    print("=" * 60)
    create_powerpoint(slides)

    print("\n" + "=" * 60)
    print("✅ All done!")
    print("=" * 60)
    print("HTML slideshow: presentation_slideshow.html")
    print("PowerPoint: Pipeline_Builder_Databricks_Presentation_HTML.pptx")


if __name__ == "__main__":
    main()
