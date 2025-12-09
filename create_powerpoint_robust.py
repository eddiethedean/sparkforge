#!/usr/bin/env python3
"""
Robust Python script to create PowerPoint presentation
with all slides from the outline.
"""

from typing import Dict, List

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError:
    print("ERROR: python-pptx not installed.")
    print("Please install it with: pip install python-pptx")
    exit(1)


# Color scheme
COLORS = {
    "bronze": RGBColor(205, 127, 50),  # #CD7F32
    "silver": RGBColor(192, 192, 192),  # #C0C0C0
    "gold": RGBColor(255, 215, 0),  # #FFD700
    "accent": RGBColor(0, 115, 230),  # Databricks blue
    "text": RGBColor(0, 0, 0),  # Black
    "background": RGBColor(255, 255, 255),  # White
    "code_bg": RGBColor(245, 245, 245),  # Light gray for code
}


def create_presentation():
    """Create the complete PowerPoint presentation."""
    print("Creating PowerPoint presentation...")

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Define all slides
    slides_data = get_all_slides_data()

    print(f"Creating {len(slides_data)} slides...")

    for i, slide_data in enumerate(slides_data, 1):
        print(f"  [{i}/{len(slides_data)}] {slide_data.get('title', 'Untitled')}")
        create_slide(prs, slide_data)

    # Save presentation
    output_file = "Pipeline_Builder_Databricks_Presentation.pptx"
    print(f"\nSaving presentation to: {output_file}")
    prs.save(output_file)
    print("✅ Presentation created successfully!")
    print(f"   File: {output_file}")
    print(f"   Total slides: {len(slides_data)}")

    return output_file


def get_all_slides_data():
    """Get all slide data - manually defined based on outline."""
    return [
        # SECTION 1: INTRODUCTION
        {
            "title": "Pipeline Builder on Databricks",
            "subtitle": "Creating and Maintaining Silver & Gold Tables\nwith Scheduled Incremental Runs",
            "type": "title",
        },
        {
            "title": "Agenda",
            "bullets": [
                "What is Pipeline Builder?",
                "Medallion Architecture Overview",
                "Building Silver & Gold Tables",
                "Incremental Processing & Scheduling",
                "Production Deployment on Databricks",
                "Real-World Example",
                "Benefits & ROI",
                "Q&A",
            ],
            "type": "bullets",
        },
        {
            "title": "The Challenge: Managing Data Pipelines at Scale",
            "left_col": [
                "200+ lines of complex Spark code",
                "Manual dependency management",
                "Scattered validation logic",
                "Difficult debugging",
                "No built-in error handling",
                "Manual schema management",
            ],
            "right_col": [
                "Data quality issues",
                "Slow time-to-market",
                "High maintenance costs",
                "Lack of visibility",
                "Incorrect business decisions",
                "Resource consumption",
            ],
            "type": "two-column",
        },
        {
            "title": "What is Pipeline Builder?",
            "subtitle": "Production-Ready Data Pipeline Framework",
            "bullets": [
                "Transforms complex Spark development into clean, maintainable code",
                "Built on Medallion Architecture (Bronze → Silver → Gold)",
                "Key Benefits:",
                "  ✓ 70% less boilerplate code",
                "  ✓ Automatic dependency management",
                "  ✓ Built-in validation",
                "  ✓ Seamless incremental processing",
            ],
            "type": "bullets",
        },
        {
            "title": "Why Databricks + Pipeline Builder?",
            "left_col": [
                "Managed Spark clusters",
                "Delta Lake (ACID transactions)",
                "Job scheduling",
                "Enterprise security",
                "Auto-scaling",
                "Time travel",
            ],
            "right_col": [
                "70% less code",
                "Automatic dependencies",
                "Built-in validation",
                "Production-ready patterns",
                "Step-by-step debugging",
                "Incremental processing",
            ],
            "bottom_text": "Together = Powerful Platform",
            "type": "two-column",
        },
        {
            "title": "What You'll Learn Today",
            "bullets": [
                "1. How to build Silver and Gold tables on Databricks",
                "2. How to schedule incremental updates automatically",
                "3. Best practices for production deployment",
                "4. Real-world examples and patterns",
                "5. How to monitor and maintain pipelines",
            ],
            "type": "bullets",
        },
        # SECTION 2: CORE CONCEPTS
        {
            "title": "What is the Medallion Architecture?",
            "content": "Data lakehouse design pattern that organizes data into three quality layers:\n\n• Bronze: Raw data (as-is, validated)\n• Silver: Cleaned and enriched data\n• Gold: Business-ready analytics",
            "type": "content",
        },
        {
            "title": "Bronze Layer: Raw Data Ingestion",
            "subtitle": '"Raw Materials Warehouse"',
            "bullets": [
                "Preserve source data exactly as received",
                "Stores in Delta Lake format",
                "Applies basic validation",
                "Tracks metadata and lineage",
                "Supports incremental processing",
                "Examples: Raw customer events, IoT sensors, transaction logs",
            ],
            "type": "bullets",
            "color": "bronze",
        },
        {
            "title": "Silver Layer: Cleaned and Enriched Data",
            "subtitle": '"Quality Control Center"',
            "bullets": [
                "Clean, standardize, and enrich data",
                "Transforms Bronze data",
                "Applies business rules",
                "Higher quality standards (95%+)",
                "Supports incremental updates",
                "Examples: Customer events with profiles, sensor data with metrics",
            ],
            "type": "bullets",
            "color": "silver",
        },
        {
            "title": "Gold Layer: Business Analytics",
            "subtitle": '"Finished Products Ready for Customers"',
            "bullets": [
                "Aggregated metrics for dashboards and reports",
                "Aggregates Silver data",
                "Creates KPIs and summaries",
                "Optimized for fast queries",
                "Highest quality standards (98%+)",
                "Examples: Daily revenue by category, customer lifetime value",
            ],
            "type": "bullets",
            "color": "gold",
        },
        {
            "title": "Data Flow: Bronze → Silver → Gold",
            "content": "Raw Data Sources\n    ↓\nBRONZE LAYER (Raw Data • Validated)\n    ↓\nSILVER LAYER (Cleaned • Enriched)\n    ↓\nGOLD LAYER (Aggregated • Analytics Ready)",
            "type": "content",
        },
        {
            "title": "Benefits of Medallion Architecture",
            "bullets": [
                "Data Quality: Progressive quality gates",
                "Reliability: Independent reprocessing",
                "Performance: Optimized storage and parallel processing",
                "Maintainability: Clear dependencies and easy debugging",
            ],
            "type": "bullets",
        },
        # SECTION 3: TECHNICAL DEEP DIVE
        {
            "title": "How Pipeline Builder Works",
            "subtitle": "Fluent API for Intuitive Pipeline Construction",
            "bullets": [
                "PipelineBuilder: Constructs pipelines",
                "PipelineRunner: Executes pipelines",
                "Validation System: Enforces data quality",
            ],
            "code": 'from pipeline_builder import PipelineBuilder\n\nbuilder = PipelineBuilder(\n    spark=spark,\n    schema="analytics"\n)',
            "type": "code",
        },
        {
            "title": "Validation System: Progressive Quality Gates",
            "content": "Bronze: 90%  →  Basic validation\nSilver: 95%  →  Business rules\nGold:   98%  →  Final quality check",
            "code": 'rules = {\n    "user_id": ["not_null"],\n    "age": ["gt", 0],\n    "status": ["in", ["active", "inactive"]]\n}',
            "type": "code",
        },
        {
            "title": "Parallel Execution: 3-5x Faster",
            "left_col": [
                "Sequential Execution:",
                "Step A: 2s",
                "Step B: 2s",
                "Step C: 2s",
                "Total: 6s",
            ],
            "right_col": [
                "Parallel Execution:",
                "Step A: 2s ┐",
                "Step B: 2s ├─ All concurrent",
                "Step C: 2s ┘",
                "Total: 2s ⚡",
            ],
            "type": "two-column",
        },
        {
            "title": "Automatic Dependency Management",
            "bullets": [
                "Detects dependencies between steps",
                "Validates all dependencies exist",
                "Orders execution correctly",
                "Prevents circular dependencies",
            ],
            "code": '# Pipeline Builder automatically knows dependencies\nbuilder.with_bronze_rules(name="events", ...)\nbuilder.add_silver_transform(\n    name="clean",\n    source_bronze="events", ...)',
            "type": "code",
        },
        {
            "title": "Code Comparison: Before vs After",
            "left_col": [
                "Before (Raw Spark):",
                "200+ lines of complex code",
                "Manual dependency management",
                "Scattered validation",
                "Difficult debugging",
            ],
            "right_col": [
                "After (Pipeline Builder):",
                "20-30 lines of clean code",
                "Automatic dependencies",
                "Centralized validation",
                "Easy debugging",
            ],
            "bottom_text": "90% Code Reduction • 87% Faster Development",
            "type": "two-column",
        },
        # SECTION 4: DATABRICKS DEPLOYMENT
        {
            "title": "Databricks Deployment Patterns",
            "left_col": [
                "Pattern 1: Single Notebook",
                "• Simple workflows",
                "• Quick prototypes",
                "• Sequential execution",
                "• Single cluster",
            ],
            "right_col": [
                "Pattern 2: Multi-Task Job",
                "• Production workflows",
                "• Complex pipelines",
                "• Different resources",
                "• Independent scaling",
            ],
            "type": "two-column",
        },
        {
            "title": "Pattern 1: Single Notebook",
            "bullets": [
                "All code in one notebook",
                "Sequential execution",
                "Simple to understand",
                "Easy to test locally",
            ],
            "code": '# Configuration\ntarget_schema = "analytics"\n\n# Build pipeline\nbuilder = PipelineBuilder(...)\nbuilder.with_bronze_rules(...)\nbuilder.add_silver_transform(...)\nbuilder.add_gold_transform(...)\n\n# Execute\npipeline = builder.to_pipeline()\nresult = pipeline.run_initial_load(...)',
            "type": "code",
        },
        {
            "title": "Pattern 2: Multi-Task Job",
            "bullets": [
                "Separate tasks with dependencies",
                "Independent scaling per task",
                "Different cluster types",
                "Production-ready",
            ],
            "content": "Task Flow:\nTask 1: Bronze Ingestion\n    ↓\nTask 2: Silver Processing\n    ↓\nTask 3: Gold Analytics",
            "type": "content",
        },
        {
            "title": "Job Configuration: Cluster & Tasks",
            "bullets": [
                "Auto-scaling: 2-8 workers",
                "Delta optimizations enabled",
                "Adaptive query execution",
                "Node type: i3.xlarge",
            ],
            "code": '{\n  "job_clusters": [{\n    "new_cluster": {\n      "spark_version": "13.3.x-scala2.12",\n      "node_type_id": "i3.xlarge",\n      "autoscale": {\n        "min_workers": 2,\n        "max_workers": 8\n      }\n    }\n  }]\n}',
            "type": "code",
        },
        {
            "title": "Scheduling: Automate Your Pipelines",
            "content": 'Daily at 2 AM: "0 0 2 * * ?"\n\nCommon Patterns:\n• Daily: "0 0 2 * * ?"\n• Hourly: "0 0 */1 * * ?"\n• Weekly: "0 0 0 * * MON"',
            "code": "# Check if initial load needed\nif tables_exist:\n    result = pipeline.run_incremental(...)\nelse:\n    result = pipeline.run_initial_load(...)",
            "type": "code",
        },
        {
            "title": "Cluster Optimization Tips",
            "left_col": [
                "For Bronze/Silver/Gold:",
                "• Larger nodes: i3.xlarge+",
                "• Auto-scaling enabled",
                "• Delta optimizations",
            ],
            "right_col": [
                "For Logging:",
                "• Smaller nodes: i3.large",
                "• Fixed size",
                "• Minimal Spark config",
            ],
            "type": "two-column",
        },
        # SECTION 5: CREATING SILVER & GOLD TABLES
        {
            "title": "Building a Complete Pipeline",
            "bullets": [
                "1. Initialize Builder",
                "2. Define Bronze Layer",
                "3. Create Silver Layer",
                "4. Create Gold Layer",
                "5. Build and Execute",
            ],
            "type": "bullets",
        },
        {
            "title": "Step 1: Initialize Builder",
            "code": 'from pipeline_builder import PipelineBuilder\nfrom pyspark.sql import SparkSession\n\n# Initialize Spark\nspark = SparkSession.builder.getOrCreate()\n\n# Create schema\ntarget_schema = "analytics"\nspark.sql(f"CREATE SCHEMA IF NOT EXISTS {target_schema}")\n\n# Initialize PipelineBuilder\nbuilder = PipelineBuilder(\n    spark=spark,\n    schema=target_schema,\n    min_bronze_rate=90.0,\n    min_silver_rate=95.0,\n    min_gold_rate=98.0\n)',
            "type": "code",
        },
        {
            "title": "Step 2: Define Bronze Layer",
            "code": '# Load source data\nsource_df = spark.read.parquet("/mnt/raw/user_events/")\n\n# Define Bronze validation rules\nbuilder.with_bronze_rules(\n    name="events",\n    rules={\n        "user_id": [F.col("user_id").isNotNull()],\n        "timestamp": [F.col("timestamp").isNotNull()],\n        "value": [F.col("value").isNotNull(), \n                  F.col("value") > 0],\n    },\n    incremental_col="timestamp"  # Enable incremental\n)',
            "bullets": [
                "✓ Validates raw data",
                "✓ Stores in Delta table: analytics.events",
                "✓ Enables incremental processing",
                "✓ Tracks metadata",
            ],
            "type": "code",
        },
        {
            "title": "Step 3: Create Silver Layer",
            "code": '# Define Silver transformation\ndef enrich_events(spark, bronze_df, prior_silvers):\n    return (\n        bronze_df\n        .withColumn("event_date", F.to_date("timestamp"))\n        .withColumn("hour", F.hour("timestamp"))\n        .filter(F.col("user_id").isNotNull())\n    )\n\n# Add Silver transform\nbuilder.add_silver_transform(\n    name="enriched_events",\n    source_bronze="events",\n    transform=enrich_events,\n    rules={\n        "user_id": [F.col("user_id").isNotNull()],\n        "event_date": [F.col("event_date").isNotNull()],\n    },\n    table_name="enriched_events"\n)',
            "type": "code",
        },
        {
            "title": "Step 4: Create Gold Layer",
            "code": '# Define Gold aggregation\ndef daily_analytics(spark, silvers):\n    return (\n        silvers["enriched_events"]\n        .groupBy("event_date")\n        .agg(\n            F.count("*").alias("total_events"),\n            F.sum("value").alias("total_revenue"),\n            F.countDistinct("user_id").alias("unique_users")\n        )\n    )\n\n# Add Gold transform\nbuilder.add_gold_transform(\n    name="daily_analytics",\n    source_silvers=["enriched_events"],\n    transform=daily_analytics,\n    rules={\n        "event_date": [F.col("event_date").isNotNull()],\n        "total_revenue": [F.col("total_revenue") >= 0],\n    },\n    table_name="daily_analytics"\n)',
            "type": "code",
        },
        {
            "title": "Step 5: Execute Pipeline",
            "code": '# Build the pipeline\npipeline = builder.to_pipeline()\n\n# Execute initial load\nresult = pipeline.run_initial_load(\n    bronze_sources={"events": source_df}\n)\n\n# Check results\nprint(f"Status: {result.status}")\nprint(f"Rows Written: {result.metrics.total_rows_written}")\nprint(f"Duration: {result.metrics.total_duration_secs:.2f}s")',
            "content": "✅ Tables Created:\n• analytics.events (Bronze)\n• analytics.enriched_events (Silver)\n• analytics.daily_analytics (Gold)",
            "type": "code",
        },
        {
            "title": "Verify Tables Created",
            "code": '# Verify Bronze table\nbronze_table = spark.table("analytics.events")\nprint(f"Bronze rows: {bronze_table.count()}")\nbronze_table.show(5)\n\n# Verify Silver table\nsilver_table = spark.table("analytics.enriched_events")\nprint(f"Silver rows: {silver_table.count()}")\nsilver_table.show(5)\n\n# Verify Gold table\ngold_table = spark.table("analytics.daily_analytics")\nprint(f"Gold rows: {gold_table.count()}")\ngold_table.show(10)',
            "type": "code",
        },
        # SECTION 6: INCREMENTAL RUNS (KEY FOCUS)
        {
            "title": "Why Incremental Processing?",
            "left_col": [
                "The Problem:",
                "• Processing all data is slow",
                "• New data arrives continuously",
                "• We only need new/changed data",
            ],
            "right_col": [
                "The Solution:",
                "• Process only new/changed data",
                "• Faster execution (minutes vs hours)",
                "• Lower compute costs (70% reduction)",
                "• More frequent updates possible",
            ],
            "type": "two-column",
        },
        {
            "title": "How Incremental Runs Work",
            "bullets": [
                "Incremental Column: Timestamp tracks data age",
                "Delta Lake Time Travel: Tracks processed data",
                "Automatic Filtering: Only processes new data",
                "",
                "Flow:",
                "1. Check last processed timestamp",
                "2. Filter new data (timestamp > last_processed)",
                "3. Process only new rows",
                "4. Update tables incrementally",
                "5. Update last processed timestamp",
            ],
            "type": "bullets",
        },
        {
            "title": "Configuring Incremental Processing",
            "code": '# Step 1: Define incremental column\nbuilder.with_bronze_rules(\n    name="events",\n    incremental_col="timestamp"  # ← Enables incremental\n)\n\n# Step 2: Run initial load (first time)\nresult = pipeline.run_initial_load(\n    bronze_sources={"events": source_df}\n)\n\n# Step 3: Run incremental updates\nnew_data_df = spark.read.parquet("/mnt/raw/user_events/")\n\nresult = pipeline.run_incremental(\n    bronze_sources={"events": new_data_df}\n)',
            "type": "code",
        },
        {
            "title": "Incremental Processing Flow",
            "content": "Initial Load (First Run)\n├─ Process ALL historical data\n├─ Create Bronze/Silver/Gold tables\n└─ Store last processed timestamp\n         ↓\nIncremental Run (Scheduled)\n├─ Load new data (since last run)\n├─ Filter: timestamp > last_processed\n├─ Process only new rows\n├─ Append to Bronze\n├─ Update Silver (incremental)\n├─ Recalculate Gold (overwrite)\n└─ Update last processed timestamp",
            "type": "content",
        },
        {
            "title": "Scheduling Incremental Updates",
            "code": '{\n  "schedule": {\n    "quartz_cron_expression": "0 0 2 * * ?",\n    "timezone_id": "America/Los_Angeles",\n    "pause_status": "UNPAUSED"\n  }\n}',
            "content": 'Option 1: Job Schedule\nDaily at 2 AM: "0 0 2 * * ?"\n\nOption 2: Programmatic\n# Check if initial load needed\ntables_exist = spark.sql(\n    f"SHOW TABLES IN {target_schema}"\n).count() > 0\n\nif tables_exist:\n    result = pipeline.run_incremental(...)\nelse:\n    result = pipeline.run_initial_load(...)',
            "type": "code",
        },
        {
            "title": "Full Refresh vs Incremental",
            "left_col": [
                "Initial Load (Full Refresh):",
                "• First time running pipeline",
                "• Schema changes",
                "• Data quality issues",
                "• Manual trigger",
            ],
            "right_col": [
                "Incremental:",
                "• Regular scheduled updates",
                "• New data arrives continuously",
                "• Standard production operation",
                "• Daily/hourly updates",
            ],
            "type": "two-column",
        },
        {
            "title": "Silver & Gold Processing Modes",
            "left_col": [
                "Silver Layer:",
                "• Incremental Append (default)",
                "  - Only new/changed data",
                "  - Faster for large datasets",
                "",
                "• Overwrite (if no incremental_col)",
                "  - Full refresh",
                "  - Ensures consistency",
            ],
            "right_col": [
                "Gold Layer:",
                "• Always Overwrite",
                "  - Aggregations recalculated",
                "  - Ensures correct totals",
                "  - Performance still good",
                "  - (Silver is incremental)",
            ],
            "type": "two-column",
        },
        {
            "title": "Monitoring Incremental Runs",
            "code": 'from pipeline_builder import LogWriter\n\nwriter = LogWriter(\n    spark=spark,\n    schema="monitoring",\n    table_name="pipeline_logs"\n)\n\n# Log initial load\nwriter.create_table(result_initial)\n\n# Log incremental runs\nwriter.append(result_inc1)\nwriter.append(result_inc2)\n\n# Query history\nlogs = spark.table("monitoring.pipeline_logs")\nrecent_runs = logs.filter(\n    "run_started_at >= current_date() - 7"\n)',
            "bullets": [
                "Tracks:",
                "✓ Run ID, timestamps, mode",
                "✓ Success/failure status",
                "✓ Rows processed/written",
                "✓ Duration, validation rates",
                "✓ Error messages (if any)",
            ],
            "type": "code",
        },
        # SECTION 7: PRODUCTION DEPLOYMENT
        {
            "title": "Production Best Practices",
            "bullets": [
                "1. Schema Management: Always create schemas explicitly",
                "2. Data Quality: Set appropriate validation thresholds",
                "3. Resource Management: Configure timeouts and retries",
                "4. Error Handling: Comprehensive try/catch blocks",
                "5. Monitoring: Use LogWriter for execution tracking",
            ],
            "type": "bullets",
        },
        {
            "title": "Error Handling & Monitoring",
            "code": 'try:\n    result = pipeline.run_incremental(...)\n    if result.status != "completed":\n        error_msg = f"Pipeline failed: {result.errors}"\n        logger.error(error_msg)\n        send_alert(error_msg)\n        raise Exception(error_msg)\nexcept Exception as e:\n    logger.error(f"Pipeline failed: {e}")\n    send_alert(f"Exception: {e}")\n    raise',
            "bullets": [
                "Monitoring:",
                "• LogWriter tracks all executions",
                "• Query execution history",
                "• Performance trends",
                "• Error tracking",
                "• Data quality metrics",
            ],
            "type": "code",
        },
        {
            "title": "Performance Optimization Tips",
            "bullets": [
                "1. Cluster Sizing: Right-size for workload",
                "2. Delta Optimizations: Auto-compact, optimize writes",
                "3. Parallel Execution: Enable for independent steps",
                "4. Data Partitioning: Partition large tables by date",
                "5. Query Optimization: Keep statistics updated",
            ],
            "type": "bullets",
        },
        {
            "title": "Security & Governance",
            "bullets": [
                "1. Access Control: Configure permissions",
                "2. Data Lineage: Automatic tracking",
                "3. Audit Logging: Comprehensive execution logs",
                "4. Schema Validation: Automatic schema checks",
            ],
            "type": "bullets",
        },
        # SECTION 8: REAL-WORLD EXAMPLE
        {
            "title": "Real-World Example: E-Commerce Analytics",
            "bullets": [
                "Business Requirements:",
                "• Bronze: Store raw customer events",
                "• Silver: Clean and enrich with profiles",
                "• Gold: Daily business metrics",
                "",
                "Pipeline:",
                "• Bronze: customer_events",
                "• Silver: enriched_events",
                "• Gold: daily_revenue",
            ],
            "type": "bullets",
        },
        {
            "title": "Complete Pipeline Code",
            "code": 'from pipeline_builder import PipelineBuilder, LogWriter\nfrom pyspark.sql import SparkSession, functions as F\n\n# Initialize\nspark = SparkSession.builder.getOrCreate()\ntarget_schema = "ecommerce_analytics"\n\n# Build pipeline\nbuilder = PipelineBuilder(spark=spark, schema=target_schema)\n\n# Bronze\nbuilder.with_bronze_rules(\n    name="customer_events",\n    rules={"user_id": ["not_null"], "price": ["gt", 0]},\n    incremental_col="timestamp"\n)\n\n# Silver\ndef enrich_events(spark, bronze_df, prior_silvers):\n    return bronze_df.withColumn("event_date", \n                                F.to_date("timestamp"))\n\nbuilder.add_silver_transform(\n    name="enriched_events",\n    source_bronze="customer_events",\n    transform=enrich_events,\n    rules={"event_date": ["not_null"]},\n    table_name="enriched_events"\n)\n\n# Gold\ndef daily_analytics(spark, silvers):\n    return silvers["enriched_events"].groupBy("event_date").agg(\n        F.sum("price").alias("total_revenue")\n    )\n\nbuilder.add_gold_transform(\n    name="daily_revenue",\n    source_silvers=["enriched_events"],\n    transform=daily_analytics,\n    rules={"total_revenue": ["gte", 0]},\n    table_name="daily_revenue"\n)\n\n# Execute\npipeline = builder.to_pipeline()\nresult = pipeline.run_incremental(...)',
            "type": "code",
        },
        {
            "title": "Results & Monitoring",
            "content": "✅ Tables Created:\n• analytics.customer_events (Bronze)\n• analytics.enriched_events (Silver)\n• analytics.daily_revenue (Gold)\n\nMonitoring:\n• Execution logged in monitoring.pipeline_logs\n• Track success rates, durations\n• Monitor data quality",
            "type": "content",
        },
        # SECTION 9: BENEFITS & ROI
        {
            "title": "Development Time Savings",
            "content": "Before:  200+ lines, 4+ hours per pipeline\nAfter:   20-30 lines, 30 minutes per pipeline\nResult:  87% faster development\n\nMetrics:\n• 90% code reduction\n• 87% faster development\n• 60% less maintenance",
            "type": "content",
        },
        {
            "title": "Performance & Cost Benefits",
            "left_col": [
                "Performance:",
                "• 3-5x faster execution (parallel)",
                "• Incremental processing",
                "  (minutes vs hours)",
                "• Better resource utilization",
            ],
            "right_col": [
                "Cost:",
                "• 70% compute cost reduction",
                "• Lower maintenance costs",
                "• Faster time-to-market",
            ],
            "type": "two-column",
        },
        {
            "title": "ROI Calculation",
            "content": "Annual Savings per Pipeline:\n• Development: $4,200\n• Maintenance: $6,000\n• Compute: $4,200\n─────────────────────\nTotal: $14,400/year\n\nFor 10 Pipelines: $144,000/year",
            "type": "content",
        },
        # SECTION 10: CONCLUSION
        {
            "title": "Key Takeaways",
            "bullets": [
                "1. Medallion Architecture provides clear data quality progression",
                "2. Pipeline Builder simplifies development (70% less code)",
                "3. Incremental processing enables efficient updates",
                "4. Databricks integration provides enterprise deployment",
                "5. Real-world examples demonstrate production patterns",
                "",
                "Next Steps:",
                "• Install: pip install pipeline_builder[pyspark]",
                "• Review examples",
                "• Build your first pipeline",
                "• Deploy to Databricks",
            ],
            "type": "bullets",
        },
        {
            "title": "Thank You!",
            "subtitle": "Questions?",
            "content": "Resources:\n• Documentation: sparkforge.readthedocs.io\n• GitHub: github.com/eddiethedean/sparkforge\n• Examples: examples/databricks/",
            "type": "title",
        },
    ]


def create_slide(prs: Presentation, slide_data: Dict):
    """Create a single slide with the given data."""
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)

    slide_type = slide_data.get("type", "content")

    if slide_type == "title":
        create_title_slide(slide, slide_data)
    elif slide_type == "bullets":
        create_bullet_slide(slide, slide_data)
    elif slide_type == "two-column":
        create_two_column_slide(slide, slide_data)
    elif slide_type == "code":
        create_code_slide(slide, slide_data)
    elif slide_type == "content":
        create_content_slide(slide, slide_data)
    else:
        create_content_slide(slide, slide_data)


def create_title_slide(slide, data):
    """Create title slide."""
    # Main title
    title = data.get("title", "")
    add_title(slide, title, font_size=44, top=Inches(2))

    # Subtitle
    if "subtitle" in data:
        add_text(slide, data["subtitle"], font_size=32, top=Inches(3.5), bold=True)

    # Footer
    if "content" in data:
        add_text(slide, data["content"], font_size=18, top=Inches(6.5))


def create_bullet_slide(slide, data):
    """Create slide with bullet points."""
    add_title(slide, data.get("title", ""), font_size=44)

    if "subtitle" in data:
        add_text(slide, data["subtitle"], font_size=32, top=Inches(1.2), bold=True)

    bullets = data.get("bullets", [])
    if bullets:
        add_bullets(slide, bullets, top=Inches(1.8 if "subtitle" in data else 1.5))

    if "code" in data:
        add_code_block(slide, data["code"], top=Inches(3))


def create_two_column_slide(slide, data):
    """Create two-column slide."""
    add_title(slide, data.get("title", ""), font_size=44)

    left_col = data.get("left_col", [])
    right_col = data.get("right_col", [])

    if left_col:
        add_bullets(
            slide,
            left_col,
            left=Inches(0.5),
            top=Inches(1.5),
            width=Inches(4.25),
            font_size=20,
        )

    if right_col:
        add_bullets(
            slide,
            right_col,
            left=Inches(5.25),
            top=Inches(1.5),
            width=Inches(4.25),
            font_size=20,
        )

    if "bottom_text" in data:
        add_text(
            slide,
            data["bottom_text"],
            font_size=28,
            top=Inches(6.5),
            bold=True,
            center=True,
        )


def create_code_slide(slide, data):
    """Create slide with code."""
    add_title(slide, data.get("title", ""), font_size=44)

    if "code" in data:
        code_text = data["code"] if isinstance(data["code"], str) else data["code"][0]
        add_code_block(slide, code_text, top=Inches(1.5))

    if "bullets" in data:
        add_bullets(slide, data["bullets"], top=Inches(5), font_size=18)

    if "content" in data:
        add_text(slide, data["content"], top=Inches(5.5), font_size=18)


def create_content_slide(slide, data):
    """Create content slide."""
    add_title(slide, data.get("title", ""), font_size=44)

    if "subtitle" in data:
        add_text(slide, data["subtitle"], font_size=32, top=Inches(1.2), bold=True)

    if "content" in data:
        add_text(slide, data["content"], top=Inches(1.8 if "subtitle" in data else 1.5))


def add_title(slide, text: str, font_size: int = 44, top: float = Inches(0.3)):
    """Add title to slide."""
    left = Inches(0.5)
    width = Inches(9)
    height = Inches(1)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP

    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = COLORS["text"]
    p.alignment = PP_ALIGN.LEFT


def add_text(
    slide,
    text: str,
    font_size: int = 24,
    top: float = Inches(1.5),
    left: float = Inches(0.5),
    width: float = Inches(9),
    bold: bool = False,
    center: bool = False,
):
    """Add text to slide."""
    height = Inches(1)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP

    lines = text.split("\n")
    for i, line in enumerate(lines):
        if line.strip():
            p = (
                text_frame.paragraphs[i]
                if i < len(text_frame.paragraphs)
                else text_frame.add_paragraph()
            )
            p.text = line.strip()
            p.font.size = Pt(font_size)
            p.font.bold = bold
            p.font.color.rgb = COLORS["text"]
            p.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT
            p.space_after = Pt(4)


def add_bullets(
    slide,
    bullets: List[str],
    font_size: int = 24,
    top: float = Inches(1.5),
    left: float = Inches(0.5),
    width: float = Inches(9),
    height: float = Inches(5.5),
):
    """Add bullet points to slide."""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP

    for i, bullet in enumerate(bullets):
        if bullet.strip():
            p = (
                text_frame.paragraphs[i]
                if i < len(text_frame.paragraphs)
                else text_frame.add_paragraph()
            )
            p.text = bullet.strip()
            p.font.size = Pt(font_size)
            p.font.color.rgb = COLORS["text"]
            p.level = 0
            p.space_after = Pt(6)


def add_code_block(slide, code: str, font_size: int = 18, top: float = Inches(1.5)):
    """Add code block to slide."""
    left = Inches(0.5)
    width = Inches(9)
    height = Inches(5.5)

    # Add background rectangle
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["code_bg"]
    shape.line.color.rgb = RGBColor(200, 200, 200)

    # Add text box on top
    textbox = slide.shapes.add_textbox(
        left + Inches(0.1), top + Inches(0.1), width - Inches(0.2), height - Inches(0.2)
    )
    text_frame = textbox.text_frame
    text_frame.word_wrap = False
    text_frame.vertical_anchor = MSO_ANCHOR.TOP

    # Use monospace font for code
    lines = code.split("\n")
    for i, line in enumerate(lines[:35]):  # Limit to 35 lines
        p = (
            text_frame.paragraphs[i]
            if i < len(text_frame.paragraphs)
            else text_frame.add_paragraph()
        )
        p.text = line
        p.font.name = "Consolas"
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLORS["text"]
        p.space_after = Pt(2)


if __name__ == "__main__":
    try:
        create_presentation()
    except Exception as e:
        print(f"ERROR: Failed to create presentation: {e}")
        import traceback

        traceback.print_exc()
