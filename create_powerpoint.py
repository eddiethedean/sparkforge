#!/usr/bin/env python3
"""
Python script to automatically create PowerPoint presentation
from the slide content document.
"""

import re
from pathlib import Path
from typing import Dict, List

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError:
    print("ERROR: python-pptx not installed.")
    print("Please install it with: pip install python-pptx")
    exit(1)


# Color scheme
COLORS = {
    "bronze": RGBColor(205, 127, 50),  # #CD7F32
    "silver": RGBColor(192, 192, 192),  # #C0C0C0
    "gold": RGBColor(255, 215, 0),  # #FFD700
    "accent": RGBColor(0, 115, 230),  # Databricks blue
    "text": RGBColor(0, 0, 0),  # Black
    "background": RGBColor(255, 255, 255),  # White
}


def parse_slide_content(content_file: str) -> List[Dict]:
    """Parse the slide content markdown file into structured data."""
    with open(content_file, encoding="utf-8") as f:
        content = f.read()

    slides = []

    # Split by slide markers
    slide_sections = re.split(r"### Slide \d+:", content)

    for section in slide_sections[1:]:  # Skip first empty section
        slide_data = {}

        # Extract slide number and title
        title_match = re.search(r"^([^\n]+)", section.strip())
        if title_match:
            slide_data["title"] = title_match.group(1).strip()

        # Extract content sections
        if "**Title" in section:
            title_match = re.search(
                r"\*\*Title[^*]+\*\*:\s*\n```\n(.+?)\n```", section, re.DOTALL
            )
            if title_match:
                slide_data["title_text"] = title_match.group(1).strip()

        # Extract content
        if "**Content" in section:
            content_match = re.search(
                r"\*\*Content[^*]+\*\*:\s*\n```\n(.+?)\n```", section, re.DOTALL
            )
            if content_match:
                slide_data["content"] = content_match.group(1).strip()

        # Extract code
        code_blocks = re.findall(r"```(?:python|json)?\n(.*?)```", section, re.DOTALL)
        if code_blocks:
            slide_data["code"] = code_blocks

        # Extract bullet points
        bullets = re.findall(r"^[•\-\*]\s+(.+)$", section, re.MULTILINE)
        if bullets:
            slide_data["bullets"] = bullets

        # Extract layout info
        if "**Layout:**" in section:
            layout_match = re.search(r"\*\*Layout:\*\*\s*(.+)", section)
            if layout_match:
                slide_data["layout"] = layout_match.group(1).strip()

        if slide_data:
            slides.append(slide_data)

    return slides


def create_presentation_from_content(content_file: str, output_file: str):
    """Create PowerPoint presentation from slide content."""
    print(f"Reading slide content from: {content_file}")

    # Read the content file
    with open(content_file, encoding="utf-8") as f:
        content = f.read()

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Parse slides from the content
    slides_data = parse_slide_content_simple(content)

    print(f"Creating {len(slides_data)} slides...")

    for i, slide_data in enumerate(slides_data, 1):
        print(f"  Creating slide {i}: {slide_data.get('title', 'Untitled')}")
        create_slide(prs, slide_data)

    # Save presentation
    print(f"\nSaving presentation to: {output_file}")
    prs.save(output_file)
    print("✅ Presentation created successfully!")
    print(f"   File: {output_file}")
    print(f"   Slides: {len(slides_data)}")


def parse_slide_content_simple(content: str) -> List[Dict]:
    """Simple parser for slide content."""
    slides = []

    # Split by slide markers - look for "### Slide X:"
    slide_pattern = r"### Slide \d+:\s*([^\n]+)"
    slide_matches = list(re.finditer(slide_pattern, content))

    if not slide_matches:
        # Try alternative pattern
        slide_pattern = r"## SECTION.*?\n\n### Slide \d+:\s*([^\n]+)"
        slide_matches = list(re.finditer(slide_pattern, content, re.MULTILINE))

    for i, match in enumerate(slide_matches):
        slide_start = match.start()
        slide_end = (
            slide_matches[i + 1].start() if i + 1 < len(slide_matches) else len(content)
        )
        slide_section = content[slide_start:slide_end]

        slide_data = {"title": match.group(1).strip()}

        # Extract title text - look for "**Title (44pt, Bold):**" followed by code block
        title_patterns = [
            r"\*\*Title[^*]+\*\*:\s*\n```\n(.+?)\n```",
            r"\*\*Title[^*]+\*\*:\s*\n```(.+?)```",
            r"Title \(44pt, Bold\):\s*\n```\n(.+?)\n```",
        ]
        for pattern in title_patterns:
            title_match = re.search(pattern, slide_section, re.DOTALL)
            if title_match:
                slide_data["title_text"] = title_match.group(1).strip()
                break

        if "title_text" not in slide_data:
            # Use the slide title as fallback
            slide_data["title_text"] = slide_data["title"]

        # Extract content - look for "**Content" or bullet points
        content_patterns = [
            r"\*\*Content[^*]+\*\*:\s*\n```\n(.+?)\n```",
            r"\*\*Content[^*]+\*\*:\s*\n```(.+?)```",
        ]
        for pattern in content_patterns:
            content_match = re.search(pattern, slide_section, re.DOTALL)
            if content_match:
                slide_data["content"] = content_match.group(1).strip()
                break

        # Extract bullet points (lines starting with •, -, or *)
        if "content" not in slide_data:
            bullets = re.findall(r"^[•\-\*]\s+(.+)$", slide_section, re.MULTILINE)
            if bullets:
                slide_data["bullets"] = bullets

        # Extract code blocks (Python, JSON, or plain code)
        code_blocks = re.findall(
            r"```(?:python|json)?\n(.*?)```", slide_section, re.DOTALL
        )
        if code_blocks:
            # Filter out very short code blocks (likely formatting artifacts)
            slide_data["code"] = [
                cb.strip() for cb in code_blocks if len(cb.strip()) > 10
            ]

        # Extract layout
        layout_match = re.search(r"\*\*Layout:\*\*\s*(.+)", slide_section)
        if layout_match:
            slide_data["layout"] = layout_match.group(1).strip()
        else:
            # Infer layout from content
            if "code" in slide_data:
                slide_data["layout"] = "code"
            elif (
                "two-column" in slide_section.lower()
                or "comparison" in slide_section.lower()
            ):
                slide_data["layout"] = "two-column"
            else:
                slide_data["layout"] = "content"

        slides.append(slide_data)

    return slides


def create_slide(prs: Presentation, slide_data: Dict):
    """Create a single slide with the given data."""
    # Use blank layout
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)

    # Add title
    title_text = slide_data.get("title_text") or slide_data.get("title", "Untitled")
    add_title(slide, title_text)

    # Add content based on layout
    layout = slide_data.get("layout", "content").lower()

    if "code" in slide_data:
        add_code_block(slide, slide_data["code"][0])
    elif "bullets" in slide_data:
        add_bullets(slide, slide_data["bullets"])
    elif "content" in slide_data:
        add_text_content(slide, slide_data["content"])
    elif "two-column" in layout or "comparison" in layout:
        # Split content for two columns
        if "bullets" in slide_data:
            mid = len(slide_data["bullets"]) // 2
            add_two_column_bullets(
                slide, slide_data["bullets"][:mid], slide_data["bullets"][mid:]
            )
    else:
        # Default: add title only or placeholder
        add_text_content(slide, slide_data.get("content", ""))


def add_title(slide, text: str, font_size: int = 44):
    """Add title to slide."""
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(1)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP

    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = COLORS["text"]
    p.alignment = PP_ALIGN.LEFT


def add_text_content(slide, text: str, font_size: int = 24):
    """Add text content to slide."""
    if not text:
        return

    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5.5)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP

    # Split by lines and add as paragraphs
    lines = text.split("\n")
    for i, line in enumerate(lines):
        if line.strip():
            p = (
                text_frame.paragraphs[i]
                if i < len(text_frame.paragraphs)
                else text_frame.add_paragraph()
            )
            p.text = line.strip()
            p.font.size = Pt(font_size)
            p.font.color.rgb = COLORS["text"]
            p.space_after = Pt(6)


def add_bullets(slide, bullets: List[str], font_size: int = 24):
    """Add bullet points to slide."""
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5.5)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP

    for i, bullet in enumerate(bullets):
        p = (
            text_frame.paragraphs[i]
            if i < len(text_frame.paragraphs)
            else text_frame.add_paragraph()
        )
        p.text = bullet.strip()
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLORS["text"]
        p.level = 0
        p.space_after = Pt(6)


def add_two_column_bullets(
    slide, left_bullets: List[str], right_bullets: List[str], font_size: int = 20
):
    """Add two columns of bullet points."""
    # Left column
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(4.25)
    height = Inches(5.5)

    left_box = slide.shapes.add_textbox(left, top, width, height)
    left_frame = left_box.text_frame
    left_frame.word_wrap = True

    for i, bullet in enumerate(left_bullets):
        p = (
            left_frame.paragraphs[i]
            if i < len(left_frame.paragraphs)
            else left_frame.add_paragraph()
        )
        p.text = bullet.strip()
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLORS["text"]
        p.space_after = Pt(6)

    # Right column
    right = Inches(5.25)
    right_box = slide.shapes.add_textbox(right, top, width, height)
    right_frame = right_box.text_frame
    right_frame.word_wrap = True

    for i, bullet in enumerate(right_bullets):
        p = (
            right_frame.paragraphs[i]
            if i < len(right_frame.paragraphs)
            else right_frame.add_paragraph()
        )
        p.text = bullet.strip()
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLORS["text"]
        p.space_after = Pt(6)


def add_code_block(slide, code: str, font_size: int = 18):
    """Add code block to slide."""
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5.5)

    # Add background rectangle
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(245, 245, 245)  # Light gray background
    shape.line.color.rgb = RGBColor(200, 200, 200)

    # Add text box on top
    textbox = slide.shapes.add_textbox(
        left + Inches(0.1), top + Inches(0.1), width - Inches(0.2), height - Inches(0.2)
    )
    text_frame = textbox.text_frame
    text_frame.word_wrap = False
    text_frame.vertical_anchor = MSO_ANCHOR.TOP

    # Use monospace font for code
    lines = code.split("\n")
    for i, line in enumerate(lines[:30]):  # Limit to 30 lines
        p = (
            text_frame.paragraphs[i]
            if i < len(text_frame.paragraphs)
            else text_frame.add_paragraph()
        )
        p.text = line
        p.font.name = "Consolas"
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLORS["text"]
        p.space_after = Pt(2)


def main():
    """Main function."""
    content_file = Path(__file__).parent / "POWERPOINT_SLIDE_CONTENT.md"
    output_file = (
        Path(__file__).parent / "Pipeline_Builder_Databricks_Presentation.pptx"
    )

    if not content_file.exists():
        print(f"ERROR: Content file not found: {content_file}")
        return

    try:
        create_presentation_from_content(str(content_file), str(output_file))
    except Exception as e:
        print(f"ERROR: Failed to create presentation: {e}")
        import traceback

        traceback.print_exc()


if __name__ == "__main__":
    main()
